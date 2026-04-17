"""
Generate PowerPoint presentation for mushroom disease classification research.
Includes: GPU redo results + RAG agent plan.

Run: python make_slides.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT   = os.path.join(BASE_DIR, "Mushroom_Disease_Classification.pptx")

# ── Colour palette ─────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT    = RGBColor(0x1E, 0xB9, 0x8A)
C_ACCENT2   = RGBColor(0xF4, 0xA2, 0x61)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHTGRAY = RGBColor(0xCC, 0xD6, 0xE0)
C_DARKGRAY  = RGBColor(0x1E, 0x2D, 0x3D)
C_GOLD      = RGBColor(0xFF, 0xD7, 0x00)
C_RED       = RGBColor(0xFF, 0x6B, 0x6B)
C_BLUE      = RGBColor(0x74, 0xB9, 0xFF)
C_PURPLE    = RGBColor(0xA2, 0x9B, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, l, t, w, h,
             size=16, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def add_multiline(slide, lines, l, t, w, h,
                  size=14, color=C_WHITE, bold=False, spacing=2,
                  align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(spacing)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb

def set_bg(slide, color):
    bg = slide.background; bg.fill.solid()
    bg.fill.fore_color.rgb = color

def header(slide, title):
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, 13.33, 0.12, fill=C_ACCENT)
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill=C_ACCENT)
    add_rect(slide, 0, 0.12, 13.33, 0.85, fill=C_DARKGRAY)
    add_text(slide, title, 0.3, 0.2, 12.5, 0.7, size=28, bold=True)
    add_rect(slide, 0.3, 0.9, len(title) * 0.155, 0.04, fill=C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, C_BG)
add_rect(s1, 0, 0, 13.33, 0.12, fill=C_ACCENT)
add_rect(s1, 0, 7.38, 13.33, 0.12, fill=C_ACCENT)
add_rect(s1, 0, 0.12, 0.5, 7.26, fill=C_DARKGRAY)
add_rect(s1, 0.08, 0.12, 0.06, 7.26, fill=C_ACCENT)

add_rect(s1, 1.0, 1.2, 11.3, 3.8, fill=C_DARKGRAY)
add_rect(s1, 1.0, 1.2, 0.08, 3.8, fill=C_ACCENT)

add_text(s1, "Mushroom Disease Classification",
         1.3, 1.4, 10.8, 1.1, size=38, bold=True)
add_text(s1, "Deep Learning & Transfer Learning on Real Farm Data",
         1.3, 2.45, 10.8, 0.65, size=22, color=C_ACCENT)
add_text(s1, "Healthy  |  Single Infected  |  Mixed Infected",
         1.3, 3.1, 10.8, 0.5, size=17, italic=True, color=C_LIGHTGRAY)
add_rect(s1, 1.3, 3.75, 10.5, 0.04, fill=C_ACCENT)
add_text(s1, "Mushroom Development Institute, Savar, Dhaka, Bangladesh  |  April 2025",
         1.3, 3.88, 10.8, 0.4, size=13, color=C_LIGHTGRAY)
add_text(s1, "761 Images   |   TensorFlow 2.x   |   Google Colab T4 GPU",
         1.3, 4.28, 10.8, 0.4, size=13, color=C_LIGHTGRAY)

# Result pills — CPU vs GPU
add_text(s1, "CPU Baseline", 1.3, 5.05, 5.5, 0.3,
         size=11, color=C_LIGHTGRAY, bold=True)
add_text(s1, "GPU Redo (Best)", 7.0, 5.05, 5.5, 0.3,
         size=11, color=C_GOLD, bold=True)

cpu_results = [("Custom CNN", "85.84%"), ("InceptionV3", "83.19%")]
gpu_results = [("ResNet50", "92.92%"), ("EfficientNetB3", "90.27%")]

for i, (name, acc) in enumerate(cpu_results):
    x = 1.3 + i * 2.9
    add_rect(s1, x, 5.4, 2.6, 0.7, fill=C_ACCENT)
    add_text(s1, f"{name}\n{acc}", x+0.1, 5.4, 2.4, 0.7,
             size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

for i, (name, acc) in enumerate(gpu_results):
    x = 7.0 + i * 2.9
    add_rect(s1, x, 5.4, 2.6, 0.7, fill=C_GOLD)
    add_text(s1, f"{name}\n{acc}", x+0.1, 5.4, 2.4, 0.7,
             size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

add_rect(s1, 6.5, 5.1, 0.04, 1.2, fill=C_ACCENT)
add_text(s1, "→  +7.08%", 6.55, 5.55, 0.8, 0.35,
         size=11, bold=True, color=C_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DATASET & METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
header(s2, "Dataset & Methodology")

add_rect(s2, 0.3, 1.1, 5.9, 5.9, fill=C_DARKGRAY)
add_rect(s2, 0.3, 1.1, 5.9, 0.45, fill=C_ACCENT)
add_text(s2, "  DATASET", 0.3, 1.1, 5.9, 0.45, size=14, bold=True, color=C_BG)
add_multiline(s2, [
    "▸  761 total images  (JPG, iPhone 11 Pro Max)",
    "▸  Location: MDI, Savar, Dhaka, Bangladesh",
    "▸  Captured: April 2025",
    "",
    "CLASS DISTRIBUTION",
    "▸  Healthy           →  299 images  (39.3%)",
    "▸  Single Infected   →  147 images  (19.3%)",
    "▸  Mixed Infected    →  315 images  (41.4%)",
    "▸  Imbalance ratio: 2.14×",
    "",
    "STRATIFIED SPLIT  (70 / 15 / 15)",
    "▸  Train  →  531  (augmented to 2,400)",
    "▸  Val    →  113 images",
    "▸  Test   →  117 images  (held out)",
    "",
    "▸  Novel: Mixed Infected class not in",
    "   any of the 4 reference papers",
], 0.5, 1.65, 5.6, 5.2, size=12, spacing=2)

add_rect(s2, 6.8, 1.1, 6.2, 5.9, fill=C_DARKGRAY)
add_rect(s2, 6.8, 1.1, 6.2, 0.45, fill=C_ACCENT2)
add_text(s2, "  METHODOLOGY", 6.8, 1.1, 6.2, 0.45, size=14, bold=True, color=C_BG)
add_multiline(s2, [
    "PHASE 1 — Data Preparation",
    "▸  Audit 761 images (0 corrupt found)",
    "▸  Stratified 70/15/15 split per class",
    "",
    "PHASE 2 — Augmentation (9 techniques)",
    "▸  Flip, Rotate ±30°, Zoom, Shear",
    "▸  Brightness, Contrast, Noise, Shift",
    "▸  531 → 2,400 training images",
    "",
    "PHASE 3 — Transfer Learning Strategy",
    "▸  Stage 1: Freeze base, train head (15 ep)",
    "▸  Stage 2: Unfreeze last 50 layers (85 ep)",
    "▸  Focal Loss + Label Smoothing (0.1)",
    "▸  Warmup + Cosine LR Decay",
    "▸  Mixed Precision (fp16) on GPU",
    "▸  Batch size: 32  |  Optimizer: Adam",
], 7.0, 1.65, 5.8, 5.2, size=12, spacing=2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CPU BASELINE vs GPU REDO (IMPROVEMENT)
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
header(s3, "Accuracy Improvement: CPU Baseline → GPU Redo")

# CPU column
add_rect(s3, 0.3, 1.1, 5.9, 0.5, fill=C_DARKGRAY)
add_text(s3, "  Phase 3 CPU Baseline", 0.3, 1.1, 5.9, 0.5,
         size=14, bold=True, color=C_LIGHTGRAY)

cpu_rows = [
    ("Custom CNN",     "85.84%", True),
    ("InceptionV3",    "83.19%", False),
    ("DenseNet201",    "73.45%", False),
    ("ResNet50",       "67.26%", False),
    ("EfficientNetB0", "44.25%", False),
    ("VGG16",          "SKIPPED", False),
]
for i, (name, acc, best) in enumerate(cpu_rows):
    y    = 1.65 + i * 0.72
    col  = C_ACCENT if best else C_DARKGRAY
    add_rect(s3, 0.3, y, 5.9, 0.68, fill=col)
    tcol = C_BG if best else C_WHITE
    add_text(s3, name, 0.5, y+0.12, 3.5, 0.45, size=13, bold=best, color=tcol)
    add_text(s3, acc,  4.2, y+0.12, 1.8, 0.45, size=13, bold=best,
             color=tcol, align=PP_ALIGN.RIGHT)
    if best:
        add_rect(s3, 0.3, y, 0.06, 0.68, fill=C_GOLD)

# Arrow
add_text(s3, "→", 6.35, 3.8, 0.7, 0.7, size=36, bold=True,
         color=C_GOLD, align=PP_ALIGN.CENTER)

# Bug callout
add_rect(s3, 6.1, 1.1, 1.2, 5.9, fill=C_BG)
add_text(s3, "Preprocessing\nbug fixed\n\nGPU training\n\nFocal Loss\n\nMore epochs\n\nBetter head",
         6.15, 1.4, 1.1, 5.4, size=10, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# GPU column
add_rect(s3, 7.4, 1.1, 5.6, 0.5, fill=C_GOLD)
add_text(s3, "  Phase 3 Redo (GPU — Colab T4)", 7.4, 1.1, 5.6, 0.5,
         size=14, bold=True, color=C_BG)

gpu_rows = [
    ("ResNet50",       "92.92%", True),
    ("VGG16",          "91.15%", False),
    ("EfficientNetB3", "90.27%", False),
    ("InceptionV3",    "89.38%", False),
    ("DenseNet201",    "87.61%", False),
    ("Custom CNN v2",  "69.91%", False),
]
for i, (name, acc, best) in enumerate(gpu_rows):
    y   = 1.65 + i * 0.72
    col = C_GOLD if best else C_DARKGRAY
    add_rect(s3, 7.4, y, 5.6, 0.68, fill=col)
    tcol = C_BG if best else C_WHITE
    add_text(s3, name, 7.6, y+0.12, 3.3, 0.45, size=13, bold=best, color=tcol)
    add_text(s3, acc,  9.8, y+0.12, 1.8, 0.45, size=13, bold=best,
             color=tcol, align=PP_ALIGN.RIGHT)

# Improvement badge
add_rect(s3, 0.3, 6.65, 12.7, 0.58, fill=C_DARKGRAY)
add_text(s3,
         "Best result improved from  85.84%  (Custom CNN, CPU)  →  92.92%  (ResNet50, GPU)    "
         "+7.08% improvement    Key fix: correct preprocessing for each architecture",
         0.5, 6.68, 12.3, 0.5, size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — KEY FINDINGS
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
header(s4, "Key Findings & Analysis")

findings = [
    (C_GOLD,    "Preprocessing Bug Fixed (+48% for EfficientNet)",
     "EfficientNetB0 got only 44.25% on CPU — it was receiving [0,1] inputs "
     "when it internally expects [0,255]. After fixing, EfficientNetB3 reached 90.27%. "
     "This was the single biggest accuracy gain in the project."),
    (C_ACCENT,  "ResNet50 Emerged as Best Model (92.92%)",
     "Surprisingly, ResNet50 outperformed DenseNet201 (ref paper: 92.5%) on this dataset. "
     "Residual connections handle the small, visually similar classes (mold textures) well. "
     "With proper GPU training, transfer learning fully matched reference paper results."),
    (C_ACCENT2, "Mixed Infected Class is Novel Contribution",
     "None of the 4 reference papers (Vidanapathirana, Shilpashree, Wongpanya, Guragain) "
     "address a mixed contamination class. Our 3-class dataset from Bangladesh fills this "
     "gap and is directly applicable to real farm scenarios."),
    (C_BLUE,    "Ensemble of Top-3 Models Expected ~94-95%",
     "Soft-voting ensemble of ResNet50 + VGG16 + EfficientNetB3 will be evaluated on the "
     "test set. Typically adds +2-3% over best single model. Full test set evaluation "
     "with confusion matrix, ROC-AUC and Grad-CAM is the next phase."),
]

for i, (col, title, body) in enumerate(findings):
    r, c = i // 2, i % 2
    x = 0.3 + c * 6.55
    y = 1.1 + r * 3.0
    add_rect(s4, x, y, 6.25, 2.75, fill=C_DARKGRAY)
    add_rect(s4, x, y, 6.25, 0.06, fill=col)
    add_rect(s4, x, y, 0.06, 2.75, fill=col)
    add_text(s4, title, x+0.2, y+0.12, 5.9, 0.45, size=14, bold=True, color=col)
    add_text(s4, body,  x+0.2, y+0.58, 5.85, 2.0,  size=11, color=C_LIGHTGRAY)

add_rect(s4, 0.3, 7.12, 12.7, 0.22, fill=C_DARKGRAY)
add_text(s4,
         "Framework: TensorFlow 2.19  |  GPU: Colab T4  |  "
         "Train: 2,400 imgs  |  Val: 113  |  Test: 117  |  "
         "Focal Loss + Warmup-Cosine LR + Mixed Precision",
         0.4, 7.13, 12.5, 0.2, size=9, color=RGBColor(0x88,0x99,0xAA),
         align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WORKFLOW DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
header(s5, "End-to-End Research Workflow")

steps = [
    ("1", "Raw Dataset\n761 Images",        C_ACCENT),
    ("2", "Preprocessing\n224×224 / RGB",   C_BLUE),
    ("3", "Augmentation\n2,400 imgs",        C_ACCENT2),
    ("4", "GPU Training\n6 Architectures",   C_PURPLE),
    ("5", "Evaluation\nF1 / Grad-CAM",       C_RED),
]
bw, gap, sx, yt = 2.0, 0.46, 0.35, 1.25
for i, (num, label, col) in enumerate(steps):
    x = sx + i * (bw + gap)
    add_rect(s5, x+0.06, yt+0.06, bw, 1.3, fill=RGBColor(0x05,0x0E,0x18))
    add_rect(s5, x, yt, bw, 1.3, fill=col)
    add_rect(s5, x+0.05, yt+0.05, 0.36, 0.36, fill=C_BG)
    add_text(s5, num, x+0.05, yt+0.03, 0.36, 0.36,
             size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s5, label, x+0.05, yt+0.44, bw-0.1, 0.8,
             size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s5, "▶", x+bw+0.05, yt+0.46, 0.36, 0.38,
                 size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

details = [
    "Healthy / Single /\nMixed folders\nFolder-based labels",
    "Resize + Normalize\nRGB Conversion\nPIL / OpenCV",
    "9 Augmentation\nTechniques\nKeras tf.data",
    "Transfer Learning\n2-stage fine-tune\nFocal Loss + GPU",
    "Test Set 117 imgs\nConfusion Matrix\nGrad-CAM Heatmap",
]
for i, detail in enumerate(details):
    x = sx + i * (bw + gap)
    add_rect(s5, x, yt+1.45, bw, 1.55, fill=C_DARKGRAY)
    add_rect(s5, x, yt+1.45, bw, 0.04, fill=steps[i][2])
    add_text(s5, detail, x+0.08, yt+1.58, bw-0.15, 1.3,
             size=11, color=C_LIGHTGRAY, align=PP_ALIGN.CENTER)

add_rect(s5, 0.35, 4.5, 12.6, 0.04, fill=RGBColor(0x2A,0x3F,0x55))
add_text(s5, "Models Compared:", 0.35, 4.62, 2.3, 0.4,
         size=12, bold=True, color=C_ACCENT)
model_pills = [
    ("Custom CNN v2", C_ACCENT),("EfficientNetB3", C_BLUE),
    ("DenseNet201",   C_PURPLE), ("InceptionV3",    C_ACCENT2),
    ("ResNet50 ★",    C_GOLD),   ("VGG16",          C_RED),
]
for i, (m, mc) in enumerate(model_pills):
    x = 2.75 + i * 1.77
    add_rect(s5, x, 4.63, 1.62, 0.38, fill=mc)
    add_text(s5, m, x, 4.63, 1.62, 0.38, size=10, bold=True,
             color=C_BG, align=PP_ALIGN.CENTER)

# Reference papers
add_rect(s5, 0.35, 5.2, 12.6, 1.55, fill=C_DARKGRAY)
add_rect(s5, 0.35, 5.2, 12.6, 0.04, fill=C_ACCENT2)
add_text(s5, "Reference Papers:", 0.55, 5.3, 2.3, 0.35,
         size=12, bold=True, color=C_ACCENT2)
refs = [
    "Vidanapathirana et al. (2023) — VGG16 / ResNet50 / InceptionV3",
    "Shilpashree et al. (2025) — YOLOv5 detection, 90% acc",
    "Wongpanya et al. (2025) — DenseNet201, 92.5% acc",
    "Guragain et al. (2024) — Custom CNN + DenseNet, 98.33% acc",
]
for i, ref in enumerate(refs):
    cx = 0.55 if i < 2 else 6.9
    ry = 5.3 + (i % 2) * 0.52
    add_text(s5, f"▸  {ref}", cx, ry+0.4, 5.9, 0.45, size=11, color=C_LIGHTGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RAG AGENT PLAN
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
header(s6, "Future Work: RAG-Based Mushroom Disease Q&A Agent")

# Description
add_rect(s6, 0.3, 1.0, 12.7, 0.55, fill=C_DARKGRAY)
add_rect(s6, 0.3, 1.0, 0.06, 0.55, fill=C_ACCENT)
add_text(s6,
         "A Retrieval-Augmented Generation (RAG) agent that answers any question about "
         "mushroom diseases, recommends treatments, and integrates with the CNN classifier.",
         0.5, 1.05, 12.3, 0.45, size=13, color=C_LIGHTGRAY)

# Architecture flow (top row)
arch_steps = [
    ("User Query\nor Image",       C_ACCENT2),
    ("CNN Classifier\n(ResNet50)", C_GOLD),
    ("Vector DB\n(ChromaDB)",      C_BLUE),
    ("LLM\n(Claude API)",          C_PURPLE),
    ("Answer +\nCitations",        C_ACCENT),
]
bw2 = 2.0
for i, (label, col) in enumerate(arch_steps):
    x = 0.35 + i * (bw2 + 0.3)
    add_rect(s6, x+0.05, 1.72, bw2, 1.0, fill=RGBColor(0x05,0x0E,0x18))
    add_rect(s6, x, 1.65, bw2, 1.0, fill=col)
    add_text(s6, label, x, 1.65, bw2, 1.0, size=12, bold=True,
             color=C_BG, align=PP_ALIGN.CENTER)
    if i < 4:
        ax = x + bw2 + 0.02
        add_text(s6, "→", ax, 1.9, 0.28, 0.5, size=20, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

# Knowledge base sources
add_rect(s6, 0.3, 2.82, 12.7, 0.38, fill=C_DARKGRAY)
add_text(s6,
         "Knowledge Base:  4 reference papers (PDFs)  +  Disease descriptions  +  "
         "Treatment protocols  +  Phase 3/4 results  +  Farm management guides",
         0.5, 2.86, 12.3, 0.3, size=11, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# Implementation phases (bottom)
phases = [
    (C_ACCENT,  "Phase 1\nKnowledge Base",
     "Collect & parse:\n4 reference PDFs\nDisease descriptions\nTreatment guides\nOur dataset results"),
    (C_BLUE,    "Phase 2\nEmbed & Index",
     "Chunk documents\nGenerate embeddings\n(sentence-transformers)\nStore in ChromaDB\nBuild FAISS index"),
    (C_PURPLE,  "Phase 3\nRAG Pipeline",
     "Query → Retrieve\nTop-K similar chunks\nBuild LLM prompt\nCall Claude API\nReturn grounded answer"),
    (C_ACCENT2, "Phase 4\nStreamlit App",
     "Chat interface\nImage upload → CNN\nGrad-CAM overlay\nCitation display\nExport Q&A report"),
]
for i, (col, title, body) in enumerate(phases):
    x = 0.3 + i * 3.27
    add_rect(s6, x, 3.3, 3.0, 3.6, fill=C_DARKGRAY)
    add_rect(s6, x, 3.3, 3.0, 0.06, fill=col)
    add_rect(s6, x, 3.3, 0.06, 3.6, fill=col)
    add_text(s6, title, x+0.2, 3.35, 2.7, 0.55, size=13, bold=True, color=col)
    add_multiline(s6, body.split("\n"), x+0.2, 3.95, 2.7, 2.8, size=11,
                  color=C_LIGHTGRAY, spacing=3)

# Tech stack
add_rect(s6, 0.3, 7.05, 12.7, 0.3, fill=C_DARKGRAY)
add_text(s6,
         "Tech Stack:  Python  |  LangChain / LlamaIndex  |  ChromaDB / FAISS  |  "
         "sentence-transformers  |  Claude API (claude-sonnet-4-6)  |  Streamlit",
         0.4, 7.07, 12.5, 0.25, size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CONCLUSION & NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
header(s7, "Conclusion & Next Steps")

add_rect(s7, 0.3, 1.1, 5.9, 5.7, fill=C_DARKGRAY)
add_rect(s7, 0.3, 1.1, 5.9, 0.45, fill=C_ACCENT)
add_text(s7, "  CONCLUSION", 0.3, 1.1, 5.9, 0.45, size=14, bold=True, color=C_BG)
add_multiline(s7, [
    "▸  3-class mushroom disease classification",
    "   successfully built end-to-end",
    "",
    "▸  CPU baseline: Custom CNN  85.84%",
    "▸  GPU redo:     ResNet50    92.92%",
    "▸  Improvement:  +7.08%",
    "",
    "▸  Root cause identified: preprocessing",
    "   mismatch caused EfficientNet to fail",
    "   (44% → 90% after fix)",
    "",
    "▸  Mixed Infected class is novel —",
    "   not addressed in any reference paper",
    "",
    "▸  Real farm dataset from Bangladesh",
    "   (761 images, MDI Savar, April 2025)",
    "",
    "▸  Full test set evaluation pending",
    "   (117 images, Phase 4)",
], 0.5, 1.65, 5.6, 5.0, size=12, spacing=2)

add_rect(s7, 6.8, 1.1, 6.2, 5.7, fill=C_DARKGRAY)
add_rect(s7, 6.8, 1.1, 6.2, 0.45, fill=C_GOLD)
add_text(s7, "  NEXT STEPS", 6.8, 1.1, 6.2, 0.45, size=14, bold=True, color=C_BG)

next_items = [
    (C_GOLD,    "Phase 4 — Test Set Evaluation",
     "Confusion matrix, ROC-AUC, per-class F1\nEnsemble top-3 models (~94-95% expected)"),
    (C_ACCENT,  "Grad-CAM Visualisation",
     "Heatmaps showing which regions the model\nfocuses on — validates mold detection"),
    (C_BLUE,    "Streamlit Web App (Phase 5)",
     "Upload image → prediction + confidence\nGrad-CAM overlay for farmer use"),
    (C_PURPLE,  "RAG Q&A Agent",
     "LangChain + ChromaDB + Claude API\nAnswer any mushroom disease question"),
    (C_ACCENT2, "Paper Writing & Submission",
     "Document 6-model comparison, novel\nMixed class, GPU improvement story"),
]
for i, (col, title, body) in enumerate(next_items):
    y = 1.65 + i * 1.0
    add_rect(s7, 6.85, y, 0.06, 0.85, fill=col)
    add_text(s7, title, 7.05, y+0.02, 5.7, 0.35, size=13, bold=True, color=col)
    add_text(s7, body,  7.05, y+0.38, 5.7, 0.55, size=11, color=C_LIGHTGRAY)

add_rect(s7, 0.3, 6.9, 12.7, 0.48, fill=C_DARKGRAY)
add_rect(s7, 0.3, 6.9, 12.7, 0.05, fill=C_ACCENT)
add_text(s7,
         "Best Model: ResNet50  92.92%  |  Dataset: MDI Savar, Dhaka 2025  |  "
         "TensorFlow 2.19  |  Colab T4 GPU  |  6 Models Trained",
         0.4, 6.97, 12.5, 0.35, size=12, bold=True,
         color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: 7 total")
print(f"  1. Title (CPU vs GPU comparison)")
print(f"  2. Dataset & Methodology")
print(f"  3. CPU Baseline → GPU Redo improvement")
print(f"  4. Key Findings & Analysis")
print(f"  5. End-to-End Workflow")
print(f"  6. RAG Agent Plan")
print(f"  7. Conclusion & Next Steps")
